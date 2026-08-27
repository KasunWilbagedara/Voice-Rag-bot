import io
import logging
from typing import Optional

from pypdf import PdfReader
from docx import Document
from google import genai
from google.genai import types

from backend.config import get_api_key, is_gemini_key

logger = logging.getLogger("voicerag.document_parser")


def extract_excel_text(file_bytes: bytes) -> str:
    """Parses Excel .xlsx / .xls files and formats workbooks/tables into clean Markdown tables."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
        markdown_sections = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows = list(sheet.iter_rows(values_only=True))
            if not rows:
                continue

            non_empty_rows = [[str(cell).strip() if cell is not None else "" for cell in row] for row in rows]
            non_empty_rows = [r for r in non_empty_rows if any(cell for cell in r)]

            if not non_empty_rows:
                continue

            markdown_sections.append(f"### Sheet: {sheet_name}\n")
            header = non_empty_rows[0]
            header_str = " | ".join(header)
            separator_str = " | ".join(["---"] * len(header))
            markdown_sections.append(f"| {header_str} |")
            markdown_sections.append(f"| {separator_str} |")

            for row in non_empty_rows[1:]:
                row_str = " | ".join(row)
                markdown_sections.append(f"| {row_str} |")
            markdown_sections.append("\n")

        return "\n".join(markdown_sections).strip()
    except Exception as e:
        logger.warning(f"openpyxl failed, attempting pandas fallback: {e}")
        try:
            import pandas as pd
            excel_file = pd.ExcelFile(io.BytesIO(file_bytes))
            sections = []
            for sheet in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet)
                if not df.empty:
                    sections.append(f"### Sheet: {sheet}\n" + df.to_markdown(index=False))
            return "\n\n".join(sections).strip()
        except Exception as err:
            raise ValueError(f"Failed to parse Excel file: {err}")


def extract_pptx_text(file_bytes: bytes) -> str:
    """Parses PowerPoint .pptx presentations slide by slide extracting titles, text, tables, and notes."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        slide_texts = []

        for i, slide in enumerate(prs.slides):
            title_text = f"Slide {i + 1}"
            if slide.shapes.title and slide.shapes.title.text.strip():
                title_text += f": {slide.shapes.title.text.strip()}"
            
            slide_content = [f"### {title_text}"]
            
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue

                if shape.has_text_frame:
                    paragraphs_text = []
                    for paragraph in shape.text_frame.paragraphs:
                        t = paragraph.text.strip()
                        if t and t not in paragraphs_text:
                            paragraphs_text.append(t)
                    if paragraphs_text:
                        slide_content.append("\n".join(paragraphs_text))

                elif shape.has_table:
                    table = shape.table
                    table_rows = []
                    for row in table.rows:
                        row_text = [cell.text.strip() for cell in row.cells]
                        table_rows.append(" | ".join(row_text))
                    if table_rows:
                        slide_content.append("\n".join(table_rows))

            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_content.append(f"Notes: {notes}")

            slide_texts.append("\n".join(slide_content))

        return "\n\n".join(slide_texts).strip()
    except Exception as e:
        raise ValueError(f"Failed to parse PowerPoint presentation: {e}")


def extract_multimodal_vision_text(file_bytes: bytes, mime_type: str, custom_api_key: Optional[str] = None) -> str:
    """Uses Google GenAI SDK Multimodal Vision / OCR to extract text from images and scanned documents."""
    try:
        api_key = get_api_key(custom_api_key)
        
        if is_gemini_key(api_key):
            client = genai.Client(api_key=api_key)
            prompt = (
                "Examine this document/image thoroughly. Extract ALL visible text, tables, numbers, headers, and descriptions. "
                "Format tables as clean Markdown tables. Preserve exact spellings, titles, names, and structural headings. "
                "Output ONLY the extracted content."
            )
            image_part = types.Part.from_bytes(
                data=file_bytes,
                mime_type=mime_type if mime_type else "image/png",
            )

            models_to_try = ["gemini-2.0-flash", "gemini-1.5-flash"]
            last_err = None
            for m_name in models_to_try:
                try:
                    res = client.models.generate_content(
                        model=m_name,
                        contents=[image_part, prompt],
                    )
                    if res and res.text and res.text.strip():
                        return res.text.strip()
                except Exception as e:
                    last_err = e

            logger.warning(f"Google GenAI Multimodal OCR failed: {last_err}")
    except Exception as e:
        logger.warning(f"Multimodal OCR vision processing error: {e}")
    
    return ""


def parse_document(
    file_bytes: bytes,
    filename: str,
    mime_type: str = "",
    custom_api_key: Optional[str] = None,
) -> str:
    """Universal Document Parser handling PDF, DOCX, XLSX, PPTX, Images, CSV, JSON, HTML, Markdown, and Code."""
    ext = filename.split(".")[-1].lower() if "." in filename else ""
    mime_lower = mime_type.lower()

    if ext == "pdf" or "pdf" in mime_lower:
        text = ""
        try:
            reader = PdfReader(io.BytesIO(file_bytes))
            text_parts = []
            for i, page in enumerate(reader.pages):
                t = page.extract_text()
                if t and t.strip():
                    text_parts.append(f"--- Page {i + 1} ---\n{t.strip()}")
            text = "\n\n".join(text_parts).strip()
        except Exception as e:
            logger.warning(f"pypdf extraction issue: {e}")

        if len(text) < 50:
            vision_text = extract_multimodal_vision_text(file_bytes, "application/pdf", custom_api_key)
            if vision_text and len(vision_text) > len(text):
                return vision_text

        return text if text else "No text extracted from PDF."

    if ext in ["docx", "doc"] or "officedocument.wordprocessingml" in mime_lower:
        try:
            doc = Document(io.BytesIO(file_bytes))
            parts = []
            for p in doc.paragraphs:
                if p.text.strip():
                    if p.style and "Heading" in p.style.name:
                        parts.append(f"## {p.text.strip()}")
                    else:
                        parts.append(p.text.strip())
            
            for table in doc.tables:
                table_rows = []
                for row in table.rows:
                    row_cells = [c.text.strip() for c in row.cells]
                    table_rows.append(" | ".join(row_cells))
                if table_rows:
                    parts.append("\n".join(table_rows))

            return "\n\n".join(parts).strip()
        except Exception as e:
            raise ValueError(f"Failed to parse Word document: {e}")

    if ext in ["xlsx", "xls", "csv"] or "excel" in mime_lower or "spreadsheet" in mime_lower:
        if ext == "csv":
            try:
                return file_bytes.decode("utf-8")
            except UnicodeDecodeError:
                return file_bytes.decode("latin-1", errors="ignore")
        return extract_excel_text(file_bytes)

    if ext in ["pptx", "ppt"] or "presentation" in mime_lower:
        return extract_pptx_text(file_bytes)

    if ext in ["png", "jpg", "jpeg", "webp", "bmp", "tiff"] or "image/" in mime_lower:
        image_mime = mime_type if "image/" in mime_lower else f"image/{ext if ext != 'jpg' else 'jpeg'}"
        vision_text = extract_multimodal_vision_text(file_bytes, image_mime, custom_api_key)
        if vision_text:
            return vision_text
        raise ValueError(f"Could not extract text/OCR from image {filename}")

    if ext in ["txt", "md", "markdown", "json", "csv", "html", "xml", "py", "js", "ts", "sql", "java", "c", "cpp"]:
        try:
            return file_bytes.decode("utf-8")
        except UnicodeDecodeError:
            return file_bytes.decode("latin-1", errors="ignore")

    try:
        return file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        return file_bytes.decode("latin-1", errors="ignore")
